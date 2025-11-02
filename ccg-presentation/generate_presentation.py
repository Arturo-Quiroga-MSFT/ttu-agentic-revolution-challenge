#!/usr/bin/env python3
"""
Generate PowerPoint presentation from slide content
Uses python-pptx library for full control over formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Azure brand colors
AZURE_BLUE = RGBColor(0, 120, 212)  # #0078D4
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 49, 48)
LIGHT_GRAY = RGBColor(243, 242, 241)
GREEN = RGBColor(16, 124, 16)
RED = RGBColor(196, 49, 75)

def create_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Contoso Consulting Group\nMulti-Agent Timesheet Assistant"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = AZURE_BLUE
    
    subtitle.text = "Production-Ready Agentic AI with Microsoft Agent Framework\n\nArturo Quiroga\nCloud Solution Architect - Data & AI, Microsoft\n\nTTU Agentic Revolution Challenge\nNovember 14, 2025"
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Welcome slide - set the stage for production-ready AI solution."

def create_problem_slide(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    title = slide.shapes.title
    title.text = "The Business Problem"
    
    # Left text box
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(3))
    tf = left.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Lost Revenue:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    for text in [
        "Consultants lose 10-15% of billable time annually",
        "Travel time most commonly missed",
        "Client meetings forgotten in timesheets",
        "Manual reconciliation wastes 15-20 min/week"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(14)
        p.level = 0
    
    # Right text box
    right = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4.5), Inches(3))
    tf = right.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Business Impact:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    for text in [
        "💰 $1,000/week lost per consultant",
        "⏱️ 15-20 minutes/week wasted",
        "⚠️ Compliance risks",
        "📉 Billing inconsistencies"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.level = 0
    
    # Bottom text
    bottom = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.8))
    tf = bottom.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Solution: Three-phase evolution: Single-Agent Demo → Multi-Agent Development → Production System"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = AZURE_BLUE
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = 'Start with relatable question: "How many of you forget to log travel time?" Emphasize this is real money being lost.'

def create_evolution_slide(prs):
    """Slide 3: Solution Evolution"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Three Phases of Innovation"
    
    # Three columns
    col_width = Inches(3)
    col_height = Inches(3.5)
    y_pos = Inches(1.8)
    
    columns = [
        {
            "x": Inches(0.5),
            "title": "Phase 1: Single-Agent POC",
            "icon": "🤖",
            "subtitle": "Proof of Concept",
            "bullets": [
                "Microsoft Agent Framework",
                "Basic calendar + timesheet analysis",
                "Read-only suggestions",
                "Status: ✅ Validated"
            ]
        },
        {
            "x": Inches(3.7),
            "title": "Phase 2: Multi-Agent Dev",
            "icon": "🔷",
            "subtitle": "Specialized Architecture",
            "bullets": [
                "4 specialized agents",
                "Parallel execution (2x faster)",
                "Domain expertise",
                "Status: ✅ Performance Proven"
            ]
        },
        {
            "x": Inches(6.9),
            "title": "Phase 3: Production ⭐",
            "icon": "🚀",
            "subtitle": "Enterprise Ready",
            "bullets": [
                "5 agents (+ Approval)",
                "Write capabilities",
                "Complete audit trail",
                "Azure deployment",
                "Status: ✅ Production Ready"
            ]
        }
    ]
    
    for col in columns:
        box = slide.shapes.add_textbox(col["x"], y_pos, col_width, col_height)
        tf = box.text_frame
        tf.word_wrap = True
        
        # Title with icon
        p = tf.paragraphs[0]
        p.text = f"{col['icon']} {col['subtitle']}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Bullets
        for bullet in col["bullets"]:
            p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(12)
            p.level = 0
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = '"Each phase built on the previous. We didn\'t just build a demo - we built a production system."'

def create_architecture_slide(prs):
    """Slide 4: Multi-Agent Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-Agent Architecture - How It Works"
    
    # Placeholder for diagram
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2.5))
    tf = diagram_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[INSERT DIAGRAM: Architecture diagram showing Orchestrator coordinating 5 agents]\n\nRefer to: diagrams/architecture.md"
    p.font.size = Pt(14)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Key benefits below
    benefits_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
    tf = benefits_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Benefits:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for benefit in [
        "⚡ Parallel Execution: 2x faster than sequential",
        "🎯 Specialized Expertise: Each agent masters one domain",
        "📈 Scalable: Add agents without breaking system",
        "✅ Approval Workflow: User controls all writes",
        "📋 Audit Trail: 100% compliance"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + benefit
        p.font.size = Pt(12)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Notice the Orchestrator at top coordinates 5 specialized agents. Calendar and Timesheet agents run in parallel - that's why it's fast."

def create_agents_slide(prs):
    """Slide 5: The Five Specialized Agents"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Agent Specialization & Roles"
    
    agents = [
        ("📅 Calendar Agent", "Calendar event analysis", "Travel time, client meetings, billable events", "get_calendar_events()"),
        ("📝 Timesheet Agent", "Timesheet validation", "Logged entries, gaps, missing time", "get_timesheet_entries()"),
        ("💡 Suggestion Agent", "Recommendation synthesis", "Missing entries, cross-references data", "suggest_timesheet_entry()"),
        ("✅ Approval Agent ⭐", "Approval workflow", "Write operations, rejections, audit logging", "add_timesheet_entry()..."),
        ("💰 Revenue Agent", "Financial impact", "Lost revenue, ROI, firm-wide projections", "calculate_revenue_impact()")
    ]
    
    y_start = Inches(1.8)
    row_height = Inches(0.8)
    
    for idx, (name, expert, identifies, tool) in enumerate(agents):
        y_pos = y_start + (idx * row_height)
        
        # Agent name
        name_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(9), Inches(0.7))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{name}  |  Expert in: {expert}  |  Tool: {tool}"
        p.font.size = Pt(12)
        if "⭐" in name:
            p.font.bold = True
            p.font.color.rgb = AZURE_BLUE
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Each agent is like a specialist on your team. The Approval Agent is production-only - it's the gatekeeper for all writes."

def create_workflow_slide(prs):
    """Slide 6: Workflow Sequence"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Production Workflow - From Analysis to Approval"
    
    # Placeholder for diagram
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = diagram_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[INSERT DIAGRAM: Workflow sequence diagram]\n\nRefer to: diagrams/workflow.md"
    p.font.size = Pt(14)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Steps
    steps_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    tf = steps_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Workflow Steps:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    steps = [
        "1. Analysis Phase: Calendar + Timesheet agents (parallel)",
        "2. Synthesis Phase: Suggestion agent proposes entries",
        "3. Review Phase: User sees suggestions with rationale",
        "4. Approval Phase: User approves or rejects ⭐",
        "5. Write Phase: Approved entries written to timesheet ⭐",
        "6. Audit Phase: All actions logged for compliance ⭐"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(12)
        if "⭐" in step:
            p.font.color.rgb = AZURE_BLUE
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Notice steps 4-6 are new in production. User approval required before any write. Everything is audited."

def create_demo_slide(prs):
    """Slide 7: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Live Demonstration"
    
    # Left side - scenario
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4))
    tf = left.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Demo Scenario"
    p.font.bold = True
    p.font.size = Pt(18)
    
    content = """
User: Sarah Johnson, Senior Consultant
Period: November 13-14, 2025

Data:
• 7 calendar events
• 2 timesheet entries (3.5 hours)
• 5 events missing (8 hours unbilled)

Missing Entries Found:
✈️ Flight to Vancouver - 2 hrs
🍽️ Working lunch - 1 hr
💬 Client Q&A - 2 hrs
✈️ Return flight - 2 hrs
📝 Workshop prep - 1 hr
"""
    
    p = tf.add_paragraph()
    p.text = content.strip()
    p.font.size = Pt(12)
    
    # Right side - results
    right = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4.5), Inches(4))
    tf = right.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "[SCREENSHOT or LIVE DEMO]"
    p.font.italic = True
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nImpact:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for item in [
        "💰 $2,000 weekly revenue recovered",
        "⚡ 5 seconds analysis time",
        "✅ All entries approved and written",
        "📋 Audit trail complete"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = '"Let me show you this live..." [Run actual demo]. If demo fails, use screenshot and walk through it.'

def create_roi_slide(prs):
    """Slide 8: Business Impact & ROI"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Measurable Business Impact"
    
    # Per Consultant
    consultant_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = consultant_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Per Consultant - Annual:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = AZURE_BLUE
    
    p = tf.add_paragraph()
    p.text = "💵 $104,000 revenue recovered  |  ⏰ 12.9 hours saved  |  📊 ROI: 200,000:1 vs AI cost"
    p.font.size = Pt(16)
    
    # Firm-Wide
    firm_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    tf = firm_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Firm-Wide (50 Consultants) - Annual:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = AZURE_BLUE
    
    p = tf.add_paragraph()
    p.text = "💰 $2.6M revenue captured  |  🚀 645 hours saved  |  📈 ROI: 4,900% in Year 1"
    p.font.size = Pt(16)
    
    # Cost breakdown
    cost_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.2))
    tf = cost_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Cost Breakdown:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "Development: $50K (one-time)  |  Azure Hosting: $1,200/year  |  AI API Calls: $520/year"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Total: $52K vs $2.6M recovered = 50x return"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These aren't theoretical numbers. $250/hr billable rate × 8 missing hours/week × 50 consultants. The ROI is undeniable."

def create_production_features_slide(prs):
    """Slide 9: Production Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "What Makes It Production-Ready"
    
    features = [
        ("Approval Workflow", [
            "✅ User approval required",
            "✅ Approve/Reject buttons",
            "✅ Rejection reason capture",
            "✅ Confirmation messages",
            "✅ No automatic writes"
        ]),
        ("Audit & Compliance", [
            "📋 Complete operation logging",
            "🕐 Timestamp tracking",
            "👤 User attribution",
            "🔒 Immutable trail",
            "📊 Compliance reports",
            "❌ No deletions (add-only)"
        ]),
        ("Deployment", [
            "🐳 Docker containerized",
            "☁️ Azure Container Apps",
            "🔧 Automated deployment",
            "📚 Complete documentation",
            "🔐 Azure Key Vault secrets"
        ])
    ]
    
    x_positions = [Inches(0.5), Inches(3.7), Inches(6.9)]
    
    for idx, (feature_title, bullets) in enumerate(features):
        box = slide.shapes.add_textbox(x_positions[idx], Inches(1.8), Inches(2.8), Inches(3.5))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = feature_title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = AZURE_BLUE
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(11)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Production-ready means more than 'it works'. It means approval workflows, audit trails, and enterprise deployment."

def create_tech_stack_slide(prs):
    """Slide 10: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Technical Architecture"
    
    # Technology stack table (simplified as text)
    table_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = table_box.text_frame
    tf.word_wrap = True
    
    stack = [
        ("Framework", "Microsoft Agent Framework", "Agent orchestration"),
        ("AI Model", "Azure OpenAI GPT-4o", "Reasoning engine"),
        ("Language", "Python 3.13+", "Application code"),
        ("UI", "Streamlit", "Web interface"),
        ("Memory", "AgentThread", "Conversation state"),
        ("Deployment", "Azure Container Apps", "Production hosting"),
        ("Registry", "Azure Container Registry", "Image storage"),
        ("Security", "Azure Key Vault", "Secrets management")
    ]
    
    p = tf.paragraphs[0]
    p.text = "Technology Stack:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for layer, tech, purpose in stack:
        p = tf.add_paragraph()
        p.text = f"{layer}: {tech} - {purpose}"
        p.font.size = Pt(12)
    
    # Key features
    features_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
    tf = features_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "• Production-grade Microsoft framework  • Azure OpenAI keeps data in your tenant  • One-command deployment  • Complete observability"
    p.font.size = Pt(12)
    p.font.bold = True
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Built on Microsoft's production frameworks - not research toys. Enterprise-ready from day one."

def create_deployment_slide(prs):
    """Slide 11: Deployment & Integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Deployment & System Integration"
    
    # Left - Deployment Options
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(3.5))
    tf = left.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Deployment Options:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    options = [
        ("Azure Container Apps (Recommended)", "$50-100/month", "Serverless, auto-scaling"),
        ("Azure App Service", "$75-150/month", "PaaS, built-in monitoring"),
        ("Azure Kubernetes Service", "$200-500/month", "Enterprise-scale, multi-region")
    ]
    
    for name, cost, desc in options:
        p = tf.add_paragraph()
        p.text = f"{name}\n{cost} - {desc}"
        p.font.size = Pt(11)
        p.space_after = Pt(10)
    
    # Right - Integration Points
    right = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4.5), Inches(3.5))
    tf = right.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Integration Points:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    integrations = [
        ("Microsoft Graph API", "Calendar events, OAuth 2.0"),
        ("ERP REST APIs", "Timesheet read/write, 2-4 weeks"),
        ("Azure AD", "SSO, RBAC permissions")
    ]
    
    for name, desc in integrations:
        p = tf.add_paragraph()
        p.text = f"{name}\n{desc}"
        p.font.size = Pt(11)
        p.space_after = Pt(10)
    
    # Timeline
    timeline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    tf = timeline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Timeline:  Demo Ready: ✅ Today  |  API Integration: 2-4 weeks  |  Pilot: 2 weeks  |  Full Rollout: 6 weeks total"
    p.font.size = Pt(12)
    p.font.bold = True
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Deployment is push-button. Integration is the only remaining step - 2-4 weeks for APIs."

def create_opensource_slide(prs):
    """Slide 12: Open Source Repository"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Open Source & Resources"
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔗 GitHub Repository:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = AZURE_BLUE
    
    p = tf.add_paragraph()
    p.text = "github.com/Arturo-Quiroga-MSFT/ttu-agentic-revolution-challenge"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nWhat's Included:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for item in [
        "✅ Complete source code (all 3 versions)",
        "✅ Deployment scripts (Docker + Azure)",
        "✅ Sample data (calendar + timesheet JSON)",
        "✅ Architecture diagrams (Mermaid)",
        "✅ Comprehensive documentation"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\nQuick Start: git clone → pip install → streamlit run"
    p.font.size = Pt(14)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Deploy to Azure: ./deploy-aca.sh"
    p.font.size = Pt(14)
    p.font.bold = True
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Everything is open source - MIT license. You can run this locally in 5 minutes, deploy to Azure in 10."

def create_lessons_slide(prs):
    """Slide 13: Lessons Learned"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Key Learnings & Best Practices"
    
    # Left - What Works
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(3.5))
    tf = left.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "What Works:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN
    
    works = [
        "✅ Multi-Agent > Single-Agent",
        "   • Specialized agents = better results",
        "   • Parallel execution = 2x speed",
        "✅ Approval Workflow is Critical",
        "   • Human-in-the-loop required",
        "   • Audit trail non-negotiable",
        "✅ Function Calling Works",
        "   • Azure OpenAI reliable",
        "   • Clear descriptions = better selection"
    ]
    
    for item in works:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
    
    # Right - Best Practices
    right = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4.5), Inches(3.5))
    tf = right.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Architecture Best Practices:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = AZURE_BLUE
    
    practices = [
        "Agent Design:",
        "• One responsibility per agent",
        "• Clear, detailed instructions",
        "• Minimal tools per agent",
        "",
        "Production Considerations:",
        "• Never commit API keys",
        "• Use Azure Key Vault",
        "• Implement RBAC",
        "• Audit all writes"
    ]
    
    for item in practices:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
    
    # Bottom - Pitfalls
    bottom = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    tf = bottom.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ Common Pitfalls: Don't make one agent do everything | Don't skip approval workflows | Don't ignore audit requirements"
    p.font.size = Pt(11)
    p.font.color.rgb = RED
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Multi-agent isn't just faster, it's fundamentally better. And approval workflows aren't optional for production."

def create_next_steps_slide(prs):
    """Slide 14: Next Steps & Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Next Steps - Let's Get Started"
    
    # Three columns
    columns = [
        ("For Contoso Consulting Group", [
            "Immediate Actions:",
            "1. ✅ Approve pilot (5 consultants)",
            "2. 🔧 Begin API integration",
            "3. 👥 Identify participants",
            "4. 📅 Schedule kickoff (Nov 18)",
            "",
            "Expected Results:",
            "• $10K recovered in 2 weeks",
            "• 5 hours saved",
            "• 90%+ satisfaction"
        ]),
        ("For Other Organizations", [
            "Get Involved:",
            "1. Clone repository",
            "2. Run locally with sample data",
            "3. Test with your workflow",
            "4. Contact for integration help",
            "",
            "We Can Help:",
            "• Architecture guidance",
            "• Integration support",
            "• Custom development",
            "• Training"
        ]),
        ("Contact & Community", [
            "📧 arturo.quiroga@microsoft.com",
            "💼 LinkedIn: Arturo Quiroga",
            "🐙 GitHub: @Arturo-Quiroga-MSFT",
            "",
            "Join the Community:",
            "• Report issues",
            "• Suggest features",
            "• Submit PRs",
            "• Share use cases"
        ])
    ]
    
    x_positions = [Inches(0.5), Inches(3.7), Inches(6.9)]
    
    for idx, (col_title, items) in enumerate(columns):
        box = slide.shapes.add_textbox(x_positions[idx], Inches(1.8), Inches(2.8), Inches(3.8))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = col_title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = AZURE_BLUE
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(10)
    
    # Bottom quote
    bottom = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.6))
    tf = bottom.text_frame
    p = tf.paragraphs[0]
    p.text = '"This is just the beginning. Agentic AI is transforming business processes. Early adopters will have competitive advantage."'
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "We've built something production-ready. The code is open. The ROI is clear. Let's work together to bring this to your organization."

def create_thankyou_slide(prs):
    """Slide 15: Thank You / Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = AZURE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    questions_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1))
    tf = questions_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📧 arturo.quiroga@microsoft.com"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "🔗 github.com/Arturo-Quiroga-MSFT/ttu-agentic-revolution-challenge"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

def main():
    """Generate the complete presentation"""
    print("🎨 Generating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create all slides
    create_title_slide(prs)
    print("✅ Slide 1: Title")
    
    create_problem_slide(prs)
    print("✅ Slide 2: Problem")
    
    create_evolution_slide(prs)
    print("✅ Slide 3: Evolution")
    
    create_architecture_slide(prs)
    print("✅ Slide 4: Architecture")
    
    create_agents_slide(prs)
    print("✅ Slide 5: Agents")
    
    create_workflow_slide(prs)
    print("✅ Slide 6: Workflow")
    
    create_demo_slide(prs)
    print("✅ Slide 7: Demo")
    
    create_roi_slide(prs)
    print("✅ Slide 8: ROI")
    
    create_production_features_slide(prs)
    print("✅ Slide 9: Production Features")
    
    create_tech_stack_slide(prs)
    print("✅ Slide 10: Tech Stack")
    
    create_deployment_slide(prs)
    print("✅ Slide 11: Deployment")
    
    create_opensource_slide(prs)
    print("✅ Slide 12: Open Source")
    
    create_lessons_slide(prs)
    print("✅ Slide 13: Lessons Learned")
    
    create_next_steps_slide(prs)
    print("✅ Slide 14: Next Steps")
    
    create_thankyou_slide(prs)
    print("✅ Slide 15: Thank You")
    
    # Save presentation
    output_file = "CCG_MultiAgent_Presentation.pptx"
    prs.save(output_file)
    
    print(f"\n✨ Presentation created successfully!")
    print(f"📄 File: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n📝 Next steps:")
    print("1. Open the presentation in PowerPoint")
    print("2. Apply your preferred theme/template")
    print("3. Add architecture and workflow diagrams (slides 4 & 6)")
    print("4. Add demo screenshot (slide 7)")
    print("5. Adjust colors and fonts to match your brand")
    print("6. Review speaker notes for each slide")

if __name__ == "__main__":
    main()
